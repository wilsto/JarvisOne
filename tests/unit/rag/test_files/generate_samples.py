"""Generate sample files for testing document processing."""

import json
from pathlib import Path
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from fpdf import FPDF
import ebooklib
from ebooklib import epub

def generate_txt_file(path: Path, content: str):
    """Generate a sample text file."""
    with open(path / "sample.txt", "w", encoding="utf-8") as f:
        f.write(content)

def generate_json_file(path: Path, content: str):
    """Generate a sample JSON file."""
    data = {
        "title": content,
        "description": "This is a test JSON file"
    }
    with open(path / "sample.json", "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)

def generate_markdown_file(path: Path, content: str):
    """Generate a sample Markdown file."""
    with open(path / "sample.md", "w", encoding="utf-8") as f:
        f.write(f"# {content}\n\nThis is a test Markdown file")

def generate_pdf_file(path: Path, content: str):
    """Generate a sample PDF file."""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=16)
    pdf.cell(200, 10, txt=content, ln=1, align="C")
    pdf.output(path / "sample.pdf")

def generate_docx_file(path: Path, content: str):
    """Generate a sample DOCX file."""
    doc = Document()
    doc.add_heading(content, 0)
    doc.add_paragraph("This is a test DOCX file")
    doc.save(path / "sample.docx")

def generate_xlsx_file(path: Path, content: str):
    """Generate a sample XLSX file."""
    wb = Workbook()
    ws = wb.active
    ws["A1"] = content
    ws["A2"] = "This is a test XLSX file"
    wb.save(path / "sample.xlsx")

def generate_pptx_file(path: Path, content: str):
    """Generate a sample PPTX file."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = content
    prs.save(path / "sample.pptx")

def generate_epub_file(path: Path, content: str):
    """Generate a sample EPUB file."""
    book = epub.EpubBook()
    
    # Set metadata
    book.set_identifier('sample123')
    book.set_title(content)
    book.set_language('en')
    
    # Add chapter
    c1 = epub.EpubHtml(title='Chapter 1', file_name='chap_01.xhtml')
    c1.content = f'<h1>{content}</h1><p>This is a test EPUB file</p>'
    
    # Add chapter to book
    book.add_item(c1)
    
    # Add navigation files
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    
    # Create spine
    book.spine = ['nav', c1]
    
    # Generate epub
    epub.write_epub(path / "sample.epub", book)

def main():
    """Generate all sample files."""
    test_files_dir = Path(__file__).parent
    
    # Text Handler files
    generate_txt_file(test_files_dir, "This is a sample text file")
    generate_json_file(test_files_dir, "Sample JSON Document")
    generate_markdown_file(test_files_dir, "Sample Markdown Document")
    
    # MarkItDown Handler files
    generate_pdf_file(test_files_dir, "Sample PDF Document")
    generate_docx_file(test_files_dir, "Sample DOCX Document")
    generate_xlsx_file(test_files_dir, "Sample XLSX Document")
    generate_pptx_file(test_files_dir, "Sample PPTX Document")
    
    # EPUB Handler files
    generate_epub_file(test_files_dir, "Sample EPUB Document")

if __name__ == "__main__":
    main()
